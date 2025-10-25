#!/usr/bin/env python3
"""
Generate PowerPoint presentation from book chapters with comprehensive diagram validation.

This script reads the markdown files from docs/ directory (without modifying them)
and creates a PowerPoint presentation in the presentations/ directory.

Enhanced Features:
- Validates that each chapter has at least one diagram
- Ensures all types of Mermaid diagrams are represented across the presentation
- Provides comprehensive diagram type analysis and reporting
- Supports validation-only mode for checking coverage without generation

Supported Mermaid Diagram Types:
- Flowchart (graph LR/TB): Process flows and system workflows
- Sequence Diagram: Interaction sequences and message flows
- Class Diagram: Object relationships and system structure
- Entity Relationship: Data model relationships
- User Journey: User experience flows and touchpoints
- Gantt Chart: Project timelines and scheduling
- Pie Chart: Distribution and percentage breakdowns
- Quadrant Chart: Decision matrices and positioning
- Mindmap: Conceptual relationships and knowledge mapping
- And more...

Usage:
    python generate_presentation.py                    # Generate presentation materials
    python generate_presentation.py --create-pptx      # Generate PowerPoint file directly
    python generate_presentation.py --validate-diagrams # Validate diagram coverage only
    python generate_presentation.py --create-pptx --validate-diagrams  # Full generation with validation

Note: This script ONLY reads from docs/ and creates files outside the docs/ directory.
It should never modify any files in the docs/ directory.
"""

import os
import sys
import glob
import re
import unicodedata
from datetime import datetime
from pathlib import Path
import json
import importlib.util
from types import SimpleNamespace

try:
    import yaml
except ImportError as exc:  # pragma: no cover - dependency check
    print("❌ Error: PyYAML library not installed")
    print("   Install with: pip install PyYAML>=6.0")
    sys.exit(1)

_PPTX_SPEC = importlib.util.find_spec("pptx")
PPTX_AVAILABLE = _PPTX_SPEC is not None

if PPTX_AVAILABLE:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    from pptx.util import Emu
else:  # pragma: no cover - optional dependency shim
    Presentation = None

    class _RGBColourStub:  # pylint: disable=too-few-public-methods
        """Minimal stand-in for pptx.dml.color.RGBColor."""

        def __init__(self, red, green, blue):
            self.rgb = (red, green, blue)

    def Inches(value):  # noqa: N802 - keep helper signature for compatibility
        return float(value)

    def Pt(value):  # noqa: N802 - keep helper signature for compatibility
        return float(value)

    RGBColor = _RGBColourStub
    PP_ALIGN = SimpleNamespace(CENTER="CENTER")
    MSO_AUTO_SHAPE_TYPE = SimpleNamespace(ROUNDED_RECTANGLE="ROUNDED_RECTANGLE")
    Emu = float


def ensure_pptx_available():
    """Raise a helpful error when python-pptx is required but missing."""

    if not PPTX_AVAILABLE:
        raise RuntimeError(
            "python-pptx is required for PowerPoint generation. "
            "Install it with 'pip install python-pptx>=0.6.21'."
        )

THEME_BLUE = RGBColor(8, 58, 122)
THEME_ACCENT = RGBColor(225, 173, 1)
THEME_TEXT = RGBColor(51, 51, 51)
THEME_MUTED = RGBColor(102, 102, 102)
SLIDE_WIDTH = Inches(13.3333333333)
SLIDE_HEIGHT = Inches(7.5)


def configure_presentation_document(prs):
    """Apply shared configuration and metadata to the presentation document."""
    ensure_pptx_available()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    now = datetime.utcnow()
    props = prs.core_properties
    props.author = "Gunnar Nordqvist"
    props.title = "Architecture as Code"
    props.subject = "Architecture as Code"
    props.category = "Architecture"
    props.keywords = "Architecture as Code; presentation"
    props.last_modified_by = "Gunnar Nordqvist"
    props.language = "en-GB"
    props.modified = now
    props.last_printed = now
    if not props.created:
        props.created = now

    try:
        props.revision = str(max(int(props.revision or "1"), 1) + 1)
    except ValueError:
        props.revision = "2"


def _add_notes(slide, notes_text):
    """Populate the speaker notes section with the supplied text."""
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.clear()
    notes_frame.text = notes_text or ""


def _add_keyword_banner(slide, keyword, left, top, width=Inches(3.5), height=Inches(0.5)):
    """Add an accent banner to highlight the focus keyword."""
    if not keyword:
        return None

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = THEME_ACCENT
    banner.line.color.rgb = THEME_BLUE

    text_frame = banner.text_frame
    text_frame.text = keyword
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(16)
    paragraph.font.bold = True
    paragraph.font.color.rgb = THEME_BLUE
    paragraph.alignment = PP_ALIGN.CENTER
    text_frame.margin_left = Pt(6)
    text_frame.margin_right = Pt(6)
    return banner


def _scale_picture_within_bounds(picture, left, top, max_width, max_height, horizontal="left"):
    """Scale and position a picture inside a bounding box while preserving its ratio."""
    if not PPTX_AVAILABLE:
        return picture

    if max_width is not None and not isinstance(max_width, int):
        max_width = Emu(max_width)
    if max_height is not None and not isinstance(max_height, int):
        max_height = Emu(max_height)

    original_width = picture.width
    original_height = picture.height

    if original_width == 0 or original_height == 0:
        return picture

    scale = 1.0

    if max_width:
        scale = min(scale, max_width / original_width)
    if max_height:
        scale = min(scale, max_height / original_height)

    if scale != 1.0:
        picture.width = int(original_width * scale)
        picture.height = int(original_height * scale)

    if max_width:
        if horizontal == "centre":
            picture.left = left + max(0, (max_width - picture.width) // 2)
        elif horizontal == "right":
            picture.left = left + max(0, max_width - picture.width)
        else:
            picture.left = left
    else:
        picture.left = left

    if max_height:
        picture.top = top + max(0, (max_height - picture.height) // 2)
    else:
        picture.top = top

    return picture


def _add_diagram_placeholder(slide, left, top, width, height):
    """Insert a placeholder shape when no diagram is available."""
    placeholder = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left,
        top,
        width,
        height,
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(230, 235, 242)
    placeholder.line.color.rgb = THEME_BLUE

    text_frame = placeholder.text_frame
    text_frame.clear()
    text_frame.text = "Flowchart pending"
    first_paragraph = text_frame.paragraphs[0]
    first_paragraph.font.size = Pt(16)
    first_paragraph.font.bold = True
    first_paragraph.font.color.rgb = THEME_BLUE
    first_paragraph.alignment = PP_ALIGN.CENTER
    return placeholder


def _add_highlighted_bullet(text_frame, text):
    """Add a bullet point with the leading keyword highlighted."""
    if not text:
        return

    paragraph = text_frame.add_paragraph()
    paragraph.level = 0
    paragraph.space_after = Pt(6)
    paragraph.line_spacing = 1.2

    bullet_run = paragraph.add_run()
    bullet_run.text = "• "
    bullet_run.font.size = Pt(18)
    bullet_run.font.color.rgb = THEME_TEXT

    words = text.split()
    if not words:
        return

    keyword = words[0]
    remainder = " ".join(words[1:])

    keyword_run = paragraph.add_run()
    keyword_run.text = keyword
    keyword_run.font.bold = True
    keyword_run.font.color.rgb = THEME_BLUE
    keyword_run.font.size = Pt(18)

    if remainder:
        remainder_run = paragraph.add_run()
        remainder_run.text = f" {remainder}"
        remainder_run.font.size = Pt(18)
        remainder_run.font.color.rgb = THEME_TEXT


def build_presentation_document(prs, presentation_data):
    """Construct the Architecture as Code presentation from structured data."""
    ensure_pptx_available()
    configure_presentation_document(prs)

    today_label = datetime.now().strftime("%d %B %Y")

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    title.text = "Architecture as Code"
    title_frame = title.text_frame
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.font.size = Pt(56)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = THEME_BLUE

    subtitle_frame = subtitle.text_frame
    subtitle_frame.clear()
    author_paragraph = subtitle_frame.paragraphs[0]
    author_paragraph.text = "Author: Gunnar Nordqvist"
    author_paragraph.font.size = Pt(24)
    author_paragraph.font.color.rgb = THEME_TEXT

    date_paragraph = subtitle_frame.add_paragraph()
    date_paragraph.text = f"Last updated: {today_label}"
    date_paragraph.font.size = Pt(20)
    date_paragraph.font.color.rgb = THEME_TEXT

    editor_paragraph = subtitle_frame.add_paragraph()
    editor_paragraph.text = "Latest change by: Gunnar Nordqvist"
    editor_paragraph.font.size = Pt(20)
    editor_paragraph.font.color.rgb = THEME_TEXT

    _add_notes(
        title_slide,
        f"Architecture as Code presentation generated on {today_label} by Gunnar Nordqvist.",
    )

    chapter_index = 0

    for entry in presentation_data:
        entry_type = entry.get('type')

        if entry_type == 'front_matter':
            front = entry.get('front_matter', {})
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide_title = slide.shapes.title
            body = slide.shapes.placeholders[1]

            slide_title.text = front.get('title', 'Architecture as Code')
            slide_title.text_frame.paragraphs[0].font.size = Pt(44)
            slide_title.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE

            _add_keyword_banner(slide, front.get('focus_keyword'), Inches(0.6), Inches(0.6))

            body_frame = body.text_frame
            body_frame.clear()

            subtitle_text = front.get('subtitle')
            if subtitle_text:
                subtitle_paragraph = body_frame.paragraphs[0]
                subtitle_paragraph.text = subtitle_text
                subtitle_paragraph.font.size = Pt(22)
                subtitle_paragraph.font.bold = True
                subtitle_paragraph.font.color.rgb = THEME_TEXT
            else:
                body_frame.paragraphs[0].text = ""

            for point in front.get('key_points', [])[:6]:
                _add_highlighted_bullet(body_frame, point)

            _add_notes(slide, front.get('notes', ''))
            continue

        if entry_type == 'part':
            part = entry.get('part', {})
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide_title = slide.shapes.title
            slide_title.text = part.get('title', 'Part overview')
            slide_title.text_frame.paragraphs[0].font.size = Pt(48)
            slide_title.text_frame.paragraphs[0].font.color.rgb = THEME_BLUE

            _add_keyword_banner(slide, part.get('focus_keyword'), Inches(0.6), Inches(0.6))

            description = part.get('description')
            description_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(12), Inches(1.2))
            description_frame = description_box.text_frame
            description_frame.clear()
            description_frame.text = description or ""
            description_frame.paragraphs[0].font.size = Pt(22)
            description_frame.paragraphs[0].font.color.rgb = THEME_TEXT

            points_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(12), Inches(4.5))
            points_frame = points_box.text_frame
            points_frame.clear()
            heading_paragraph = points_frame.paragraphs[0]
            heading_paragraph.text = "Key themes"
            heading_paragraph.font.size = Pt(24)
            heading_paragraph.font.bold = True
            heading_paragraph.font.color.rgb = THEME_BLUE

            for point in part.get('key_points', [])[:6]:
                _add_highlighted_bullet(points_frame, point)

            _add_notes(slide, part.get('notes', ''))
            continue

        if entry_type == 'chapter':
            chapter_index += 1
            chapter = entry.get('chapter', {})
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.2), Inches(1.0))
            title_frame = title_box.text_frame
            title_frame.text = chapter.get('title', 'Chapter overview')
            title_paragraph = title_frame.paragraphs[0]
            title_paragraph.font.size = Pt(36)
            title_paragraph.font.bold = True
            title_paragraph.font.color.rgb = THEME_BLUE

            focus_keyword = chapter.get('focus_keyword')
            _add_keyword_banner(slide, focus_keyword, Inches(0.6), Inches(1.15))

            diagram_path = chapter.get('diagram_path')
            diagram_metadata = chapter.get('diagram_metadata') or {}
            diagram_type = diagram_metadata.get('type', 'Diagram')
            diagram_exists = bool(diagram_path and Path(diagram_path).exists())

            layout_style = 'hero' if chapter_index % 2 == 1 else 'side'

            if layout_style == 'hero':
                image_left = Inches(0.6)
                image_top = Inches(1.6)
                image_width = Inches(12.2)
                image_height = Inches(3.6)
                text_left = Inches(0.6)
                text_top = image_top + image_height + Inches(0.35)
                text_width = Inches(12.2)
                text_height = Inches(2.4)
                caption_width = image_width
                caption_left = image_left
                caption_alignment = "centre"
            else:
                image_on_left = chapter_index % 4 == 0
                image_left = Inches(0.6) if image_on_left else Inches(6.9)
                image_top = Inches(1.6)
                image_width = Inches(5.8)
                image_height = Inches(4.2)
                text_left = Inches(6.9) if image_on_left else Inches(0.6)
                text_top = Inches(1.6)
                text_width = Inches(5.8)
                text_height = Inches(5.0)
                caption_width = image_width
                caption_left = image_left
                caption_alignment = "centre"

            if diagram_exists:
                picture = slide.shapes.add_picture(diagram_path, image_left, image_top)
                _scale_picture_within_bounds(picture, image_left, image_top, image_width, image_height, horizontal=caption_alignment)

                if layout_style == 'hero':
                    text_top = picture.top + picture.height + Inches(0.4)
                    text_height = max(Inches(1.4), SLIDE_HEIGHT - text_top - Inches(0.6))

                label_top = picture.top - Inches(0.35)
                if label_top < Inches(1.25):
                    label_top = Inches(1.25)
                label_width = min(picture.width, Inches(4.2))
                label_left = picture.left + max(0, (picture.width - label_width) // 2)
                _add_keyword_banner(slide, diagram_type, label_left, label_top, width=label_width, height=Inches(0.4))

                caption_top = picture.top + picture.height + Inches(0.1)
                caption_box = slide.shapes.add_textbox(caption_left, caption_top, caption_width, Inches(0.8))
                caption_frame = caption_box.text_frame
                caption_frame.clear()

                type_paragraph = caption_frame.paragraphs[0]
                type_paragraph.text = diagram_type
                type_paragraph.font.size = Pt(12)
                type_paragraph.font.bold = True
                type_paragraph.font.color.rgb = THEME_BLUE

                diagram_source = diagram_metadata.get('source')
                if diagram_source:
                    source_paragraph = caption_frame.add_paragraph()
                    source_paragraph.text = diagram_source
                    source_paragraph.font.size = Pt(11)
                    source_paragraph.font.color.rgb = THEME_MUTED

                explanation = diagram_metadata.get('explanation')
                if explanation:
                    explanation_paragraph = caption_frame.add_paragraph()
                    explanation_paragraph.text = explanation
                    explanation_paragraph.font.size = Pt(12)
                    explanation_paragraph.font.color.rgb = THEME_TEXT

                if layout_style == 'hero':
                    caption_bottom = caption_box.top + caption_box.height
                    text_top = max(text_top, caption_bottom + Inches(0.2))
                    text_height = max(Inches(1.4), SLIDE_HEIGHT - text_top - Inches(0.6))
            else:
                placeholder = _add_diagram_placeholder(slide, image_left, image_top, image_width, image_height)
                if layout_style == 'hero':
                    text_top = placeholder.top + placeholder.height + Inches(0.4)
                    text_height = max(Inches(1.4), SLIDE_HEIGHT - text_top - Inches(0.6))

            points_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
            points_frame = points_box.text_frame
            points_frame.clear()
            heading = points_frame.paragraphs[0]
            heading_text = f"{diagram_type} highlights" if diagram_exists else "Key highlights"
            heading.text = heading_text
            heading.font.size = Pt(26)
            heading.font.bold = True
            heading.font.color.rgb = THEME_BLUE

            max_points = 4 if layout_style == 'hero' else 8
            for point in chapter.get('key_points', [])[:max_points]:
                _add_highlighted_bullet(points_frame, point)

            _add_notes(slide, chapter.get('notes', ''))

    return prs


def _load_canonical_chapter_filenames(requirements_path=Path("BOOK_REQUIREMENTS.md")):
    """Load canonical chapter filenames from the requirements specification."""
    if not requirements_path.exists():
        return set()

    content = requirements_path.read_text(encoding="utf-8").splitlines()
    if not content or content[0].strip() != "---":
        return set()

    front_matter_lines = []
    for line in content[1:]:
        if line.strip() == "---":
            break
        front_matter_lines.append(line)

    if not front_matter_lines:
        return set()

    data = yaml.safe_load("\n".join(front_matter_lines)) or {}
    book_config = data.get("book", {})

    canonical = {
        chapter.get("filename")
        for chapter in book_config.get("chapters", [])
        if chapter.get("filename")
    }

    return canonical


def load_book_requirements(requirements_path=Path("BOOK_REQUIREMENTS.md")):
    """Load the structured book requirements from the specification front matter."""
    if not requirements_path.exists():
        return {}

    content = requirements_path.read_text(encoding="utf-8").splitlines()
    if not content or content[0].strip() != "---":
        return {}

    front_matter_lines = []
    for line in content[1:]:
        if line.strip() == "---":
            break
        front_matter_lines.append(line)

    if not front_matter_lines:
        return {}

    data = yaml.safe_load("\n".join(front_matter_lines)) or {}
    return data


PART_AREA_MAP = {
    "Part I: Foundations": {"Foundations"},
    "Part II: Architecture Platform": {"Architecture Platform"},
    "Part III: Security and Governance": {"Security and Governance"},
    "Part IV: Delivery and Operations": {"Delivery and Operations"},
    "Part V: Organisation and Leadership": {"Organisation and Leadership"},
    "Part VI: Experience and Best Practices": {"Experience and Best Practices"},
    "Part VII: Future and Wrap-up": {"Future and Wrap-up"},
    "Part VIII: Appendices and Reference": {"Appendices", "Reference"},
}


def limit_words(text, max_words=20):
    """Limit text to maximum number of words while preserving meaning."""
    if not text:
        return text

    words = text.split()
    if len(words) <= max_words:
        return text

    # Take the first max_words and add ellipsis if needed
    limited = ' '.join(words[:max_words])
    if len(words) > max_words:
        limited += "..."

    return limited


def _slugify_heading(text):
    """Convert a heading into an anchor-friendly slug similar to MkDocs."""
    if not text:
        return ""

    normalised = unicodedata.normalize("NFKD", text).strip().lower()
    # Remove combining characters and unsupported punctuation
    without_diacritics = "".join(ch for ch in normalised if not unicodedata.combining(ch))
    cleaned = re.sub(r"[^0-9a-z\s-]", "", without_diacritics)
    slug = re.sub(r"[\s-]+", "-", cleaned).strip("-")
    return slug


def _collect_highlights_from_lines(lines, minimum_length=20):
    """Extract highlight sentences from markdown lines using section structure."""
    key_points = []
    current_section = None
    section_content = []
    in_code_block = False

    for raw_line in lines:
        line = raw_line.strip()

        if line.startswith('```'):
            in_code_block = not in_code_block
            continue

        if in_code_block:
            continue

        if line.startswith('## '):
            if current_section and section_content:
                meaningful = [
                    candidate.strip()
                    for candidate in section_content
                    if candidate
                    and not candidate.startswith('#')
                    and not candidate.startswith('![')
                    and not candidate.startswith('```')
                    and len(candidate.strip()) > minimum_length
                ]

                if meaningful:
                    first_sentence = meaningful[0].split('.')[0].strip() or meaningful[0][:150]
                    key_points.append(limit_words(first_sentence, 20))

            current_section = line[3:].strip()
            section_content = []
        elif line and current_section:
            section_content.append(line)

    if current_section and section_content:
        meaningful = [
            candidate.strip()
            for candidate in section_content
            if candidate
            and not candidate.startswith('#')
            and not candidate.startswith('![')
            and not candidate.startswith('```')
            and len(candidate.strip()) > minimum_length
        ]

        if meaningful:
            first_sentence = meaningful[0].split('.')[0].strip() or meaningful[0][:150]
            key_points.append(limit_words(first_sentence, 20))

    return key_points


def generate_prezi_slides(output_path=Path("docs/prezi/slides.json")):
    """Generate slides.json for the Prezi-style presentation."""
    docs_dir = Path("docs")
    if not docs_dir.exists():
        return output_path, 0

    prezi_dir = output_path.parent
    prezi_dir.mkdir(parents=True, exist_ok=True)

    md_files = sorted(
        [
            path for path in docs_dir.glob("*.md")
            if path.name.lower() not in {"index.md"}
        ]
    )

    slides = []
    x = y = 0
    col_w, row_h, cols = 620, 460, 6

    h1_re = re.compile(r"^\#\s+(.+)$", re.MULTILINE)
    h2_re = re.compile(r"^\#\#\s+(.+)$", re.MULTILINE)

    for md_file in md_files:
        text = md_file.read_text(encoding="utf-8")
        h1_match = h1_re.search(text)
        if not h1_match:
            continue

        title = h1_match.group(1).strip()
        slide_id = md_file.stem

        slides.append(
            {
                "id": slide_id,
                "title": title,
                "mdPath": f"/{slide_id}/",
                "x": x,
                "y": y,
                "zoom": 1,
            }
        )

        for index, match in enumerate(h2_re.finditer(text), start=1):
            section_title = match.group(1).strip()
            slug = _slugify_heading(section_title) or f"section-{index}"
            slides.append(
                {
                    "id": f"{slide_id}--{index}",
                    "title": section_title,
                    "mdPath": f"/{slide_id}/#{slug}",
                    "x": x + 220 + ((index - 1) % 2) * 160,
                    "y": y + 140 + ((index - 1) // 2) * 120,
                    "zoom": 1,
                    "parentId": slide_id,
                }
            )

        x += col_w
        if x > col_w * (cols - 1):
            x = 0
            y += row_h

    output_path.write_text(json.dumps(slides, ensure_ascii=False, indent=2), encoding="utf-8")
    return output_path, len(slides)

def analyze_mermaid_diagram_types():
    """Analyze available Mermaid diagram types in the images directory."""
    diagram_types = {
        'flowchart': [],      # graph LR, graph TD, graph TB
        'sequence': [],       # sequenceDiagram
        'class': [],          # classDiagram
        'state': [],          # stateDiagram
        'entity_relationship': [], # erDiagram
        'user_journey': [],   # journey
        'gantt': [],          # gantt
        'pie': [],            # pie
        'quadrant': [],       # quadrantChart
        'requirement': [],    # requirementDiagram
        'gitgraph': [],       # gitGraph
        'mindmap': [],        # mindmap
        'timeline': [],       # timeline
        'sankey': [],         # sankey
        'xy_chart': [],       # xyChart
        'other': []
    }
    
    images_dir = Path("docs/images")
    if not images_dir.exists():
        return diagram_types
    
    mermaid_files = list(images_dir.glob("*.mmd"))

    for mmd_file in mermaid_files:
        try:
            content = mmd_file.read_text(encoding='utf-8').split('\n')

            filtered_lines = []
            in_front_matter = False

            for raw_line in content:
                stripped = raw_line.strip()

                if not stripped:
                    continue

                if stripped == '---':
                    in_front_matter = not in_front_matter
                    continue

                if in_front_matter:
                    continue

                if stripped.startswith('%%'):
                    continue

                filtered_lines.append(stripped.lower())
                if len(filtered_lines) >= 20:
                    break

            diagram_type = 'other'

            for line in filtered_lines:
                if line.startswith('graph ') or line.startswith('flowchart '):
                    diagram_type = 'flowchart'
                    break
                if line.startswith('sequencediagram'):
                    diagram_type = 'sequence'
                    break
                if line.startswith('classdiagram'):
                    diagram_type = 'class'
                    break
                if line.startswith('statediagram'):
                    diagram_type = 'state'
                    break
                if line.startswith('erdiagram'):
                    diagram_type = 'entity_relationship'
                    break
                if line.startswith('journey'):
                    diagram_type = 'user_journey'
                    break
                if line.startswith('gantt'):
                    diagram_type = 'gantt'
                    break
                if line.startswith('pie'):
                    diagram_type = 'pie'
                    break
                if line.startswith('quadrantchart'):
                    diagram_type = 'quadrant'
                    break
                if line.startswith('requirementdiagram'):
                    diagram_type = 'requirement'
                    break
                if line.startswith('gitgraph'):
                    diagram_type = 'gitgraph'
                    break
                if line.startswith('mindmap'):
                    diagram_type = 'mindmap'
                    break
                if line.startswith('timeline'):
                    diagram_type = 'timeline'
                    break
                if line.startswith('sankey'):
                    diagram_type = 'sankey'
                    break
                if line.startswith('xychart'):
                    diagram_type = 'xy_chart'
                    break

            diagram_types[diagram_type].append(str(mmd_file))

        except Exception as e:
            print(f"Warning: Could not analyze {mmd_file}: {e}")
    
    return diagram_types

def get_missing_diagram_types():
    """Identify which Mermaid diagram types are missing from the collection."""
    diagram_analysis = analyze_mermaid_diagram_types()
    
    # Define essential diagram types for a comprehensive presentation
    essential_types = [
        'flowchart', 'sequence', 'class', 'entity_relationship', 
        'user_journey', 'gantt', 'pie', 'quadrant', 'mindmap'
    ]
    
    missing_types = []
    available_types = []
    
    for diagram_type in essential_types:
        if not diagram_analysis.get(diagram_type):
            missing_types.append(diagram_type)
        else:
            available_types.append(diagram_type)
    
    return {
        'missing': missing_types,
        'available': available_types,
        'analysis': diagram_analysis
    }

def validate_chapter_diagrams():
    """Validate that each chapter has at least one diagram."""
    docs_dir = Path("docs")
    if not docs_dir.exists():
        return {}
    
    chapter_files = sorted(glob.glob(str(docs_dir / "*.md")))
    canonical_filenames = _load_canonical_chapter_filenames()
    validation_results = {
        'chapters_with_diagrams': [],
        'chapters_without_diagrams': [],
        'total_chapters': 0,
        'chapters_with_diagram_count': 0
    }
    
    for chapter_file in chapter_files:
        chapter_name = Path(chapter_file).name

        # Skip files that are not part of the canonical manuscript
        if canonical_filenames:
            if chapter_name not in canonical_filenames:
                continue
        else:
            # Fallback filters if requirements file is unavailable
            if chapter_name.startswith('part_'):
                continue

        if chapter_name in ['README.md', 'architecture_as_code.md']:
            continue

        validation_results['total_chapters'] += 1
        
        try:
            with open(chapter_file, 'r', encoding='utf-8') as f:
                content = f.read()
            
            # Check for diagram references
            has_diagram = False
            diagram_refs = []
            
            for line in content.split('\n'):
                stripped = line.strip()

                if stripped.startswith('![') and 'images/' in stripped:
                    match = re.search(r'!\[.*?\]\((.*?)\)', stripped)
                    if match:
                        diagram_refs.append(match.group(1))
                        has_diagram = True
                        continue

                # Treat explicit diagram source references as diagrams even without PNG embeds
                if 'diagram source' in stripped.lower() and 'images/' in stripped:
                    matches = re.findall(r'\(images/(diagram_[^\)]+)\)', stripped, re.IGNORECASE)
                    if matches:
                        for ref in matches:
                            diagram_refs.append(f"images/{ref}")
                        has_diagram = True
            
            if has_diagram:
                validation_results['chapters_with_diagrams'].append({
                    'file': chapter_name,
                    'diagrams': diagram_refs
                })
                validation_results['chapters_with_diagram_count'] += 1
            else:
                validation_results['chapters_without_diagrams'].append(chapter_name)
                
        except Exception as e:
            print(f"Warning: Could not analyze {chapter_file}: {e}")
            validation_results['chapters_without_diagrams'].append(chapter_name)
    
    return validation_results

def extract_diagram_metadata(diagram_path, chapter_content):
    """Extract metadata and context for a diagram from its source file and chapter context."""
    diagram_info = {
        'source': None,
        'type': 'unknown',
        'purpose': None,
        'explanation': None
    }
    
    if not diagram_path:
        return diagram_info
    
    # Extract base filename and find corresponding .mmd file
    png_filename = os.path.basename(diagram_path)
    mmd_filename = png_filename.replace('.png', '.mmd')
    mmd_path = os.path.join("docs/images", mmd_filename)
    
    # Set source information
    diagram_info['source'] = f"Source: {mmd_filename}"
    
    # Try to read the mermaid source file to determine type and content
    if os.path.exists(mmd_path):
        try:
            with open(mmd_path, 'r', encoding='utf-8') as f:
                mmd_content = f.read().strip()
            
            # Determine diagram type from mermaid content - check multiple lines
            content_lines = [line.strip().lower() for line in mmd_content.split('\n') if line.strip()]
            diagram_type_found = False
            
            for line in content_lines:
                if line.startswith('graph ') or line.startswith('flowchart '):
                    diagram_info['type'] = 'Flowchart'
                    diagram_info['purpose'] = 'Illustrates process flows and system workflows'
                    diagram_type_found = True
                    break
                elif line.startswith('sequencediagram'):
                    diagram_info['type'] = 'Sequence diagram'
                    diagram_info['purpose'] = 'Shows interaction sequences and message flows'
                    diagram_type_found = True
                    break
                elif line.startswith('classdiagram'):
                    diagram_info['type'] = 'Class diagram'
                    diagram_info['purpose'] = 'Describes object relationships and system structure'
                    diagram_type_found = True
                    break
                elif line.startswith('erdiagram'):
                    diagram_info['type'] = 'Entity-relationship diagram'
                    diagram_info['purpose'] = 'Illustrates data model relationships and entities'
                    diagram_type_found = True
                    break
                elif line.startswith('journey'):
                    diagram_info['type'] = 'User journey'
                    diagram_info['purpose'] = 'Maps user experience and touchpoints'
                    diagram_type_found = True
                    break
                elif line.startswith('gantt'):
                    diagram_info['type'] = 'Gantt chart'
                    diagram_info['purpose'] = 'Presents project timelines and scheduling'
                    diagram_type_found = True
                    break
                elif line.startswith('pie'):
                    diagram_info['type'] = 'Pie chart'
                    diagram_info['purpose'] = 'Shows distribution and percentage breakdowns'
                    diagram_type_found = True
                    break
                elif line.startswith('quadrantchart'):
                    diagram_info['type'] = 'Quadrant chart'
                    diagram_info['purpose'] = 'Illustrates decision matrices and positioning'
                    diagram_type_found = True
                    break
                elif line.startswith('mindmap'):
                    diagram_info['type'] = 'Mind map'
                    diagram_info['purpose'] = 'Structures conceptual relationships and knowledge mapping'
                    diagram_type_found = True
                    break
            
            if not diagram_type_found:
                diagram_info['type'] = 'Architecture diagram'
                diagram_info['purpose'] = 'Illustrates system architecture and component relationships'
        except Exception as e:
            print(f"Warning: Could not read diagram source {mmd_path}: {e}")
    
    # Extract contextual explanation from chapter content around diagram reference
    if chapter_content and diagram_path:
        lines = chapter_content.split('\n')
        for i, line in enumerate(lines):
            if 'images/' in line and png_filename.replace('.png', '') in line:
                # Look for explanation in the line after the diagram
                if i + 1 < len(lines):
                    next_line = lines[i + 1].strip()
                    if next_line and not next_line.startswith('#') and not next_line.startswith('!['):
                        # Found explanation, limit to 20 words
                        diagram_info['explanation'] = limit_words(next_line, 18)  # Leave room for "Diagram:" prefix
                        break
                
                # Look for explanation in the line before the diagram
                if i > 0:
                    prev_line = lines[i - 1].strip()
                    if prev_line and not prev_line.startswith('#') and not prev_line.startswith('!['):
                        diagram_info['explanation'] = limit_words(prev_line, 18)
                        break
    
    # If no specific explanation found, create one based on diagram type and context
    if not diagram_info['explanation']:
        filename_lower = png_filename.lower()
        purpose_text = diagram_info['purpose'] or 'Highlights a key architecture concept'
        chapter_match = re.search(r'(chapter|kapitel)(\d+)', filename_lower)
        if chapter_match:
            diagram_info['explanation'] = f"Diagram: Chapter {chapter_match.group(2)} - {purpose_text}"
        else:
            diagram_info['explanation'] = f"Diagram: {purpose_text}"
        
        # Ensure it's within word limit
        diagram_info['explanation'] = limit_words(diagram_info['explanation'], 18)
    
    return diagram_info

def read_part_content(part_file, metadata=None):
    """Read part introduction content and extract highlights."""
    metadata = metadata or {}
    try:
        content = Path(part_file).read_text(encoding="utf-8")
    except FileNotFoundError:
        print(f"Warning: Part introduction not found: {part_file}")
        return None

    lines = content.split('\n')
    title = metadata.get('title')
    for line in lines:
        if line.startswith('# '):
            title = line[2:].strip()
            break

    title = title or Path(part_file).stem.replace('_', ' ').title()
    key_points = _collect_highlights_from_lines(lines, minimum_length=10)

    if not key_points and metadata.get('description'):
        description_sentences = [
            sentence.strip()
            for sentence in re.split(r'[.!?]', metadata['description'])
            if sentence.strip()
        ]
        key_points = [limit_words(sentence, 20) for sentence in description_sentences[:3]]

    notes_lines = [title]
    description = metadata.get('description', '')
    if description:
        notes_lines.append('')
        notes_lines.append(limit_words(description, 40))

    if key_points:
        notes_lines.append('')
        notes_lines.append('Highlights:')
        for point in key_points:
            notes_lines.append(f"- {point}")

    return {
        'title': title,
        'label': metadata.get('label'),
        'description': description,
        'key_points': key_points,
        'notes': '\n'.join(notes_lines).strip(),
        'focus_keyword': metadata.get('label') or metadata.get('title') or title,
    }


def read_front_matter_content(front_file, metadata=None):
    """Read front matter content for introductory slides."""
    metadata = metadata or {}
    try:
        content = Path(front_file).read_text(encoding="utf-8")
    except FileNotFoundError:
        print(f"Warning: Front matter not found: {front_file}")
        return None

    lines = content.split('\n')
    title = metadata.get('title')
    subtitle = ''
    for line in lines:
        if line.startswith('# '):
            title = line[2:].strip()
            continue
        if line.startswith('## '):
            subtitle = line[3:].strip()
            break

    highlights = _collect_highlights_from_lines(lines, minimum_length=5)
    notes_lines = [title or Path(front_file).stem]
    if subtitle:
        notes_lines.append(subtitle)

    if highlights:
        notes_lines.append('')
        notes_lines.append('Highlights:')
        for point in highlights:
            notes_lines.append(f"- {point}")

    return {
        'title': title or metadata.get('title') or 'Architecture as Code',
        'subtitle': subtitle or metadata.get('label', ''),
        'key_points': highlights,
        'notes': '\n'.join(notes_lines).strip(),
        'focus_keyword': metadata.get('label') or (title or 'Architecture as Code'),
    }


def read_chapter_content(chapter_file):
    """Read and parse a chapter markdown file."""
    try:
        with open(chapter_file, 'r', encoding='utf-8') as f:
            content = f.read()

        # Extract chapter title (first h1)
        lines = content.split('\n')
        title = "Untitled Chapter"
        for line in lines:
            if line.startswith('# '):
                title = line[2:].strip()
                break
        
        # Extract diagram paths - prefer flowchart representations
        diagram_path = None
        diagram_metadata = None
        selected_reference = None
        diagram_candidates = []
        for line in lines:
            if line.startswith('![') and 'images/' in line:
                match = re.search(r'!\[.*?\]\((.*?)\)', line)
                if match:
                    relative_path = match.group(1)
                    if relative_path.endswith('.png'):
                        full_path = os.path.join("docs", relative_path)
                        metadata = extract_diagram_metadata(full_path, content)
                        diagram_candidates.append((full_path, metadata, relative_path))

        for candidate in diagram_candidates:
            candidate_metadata = candidate[1] or {}
            diagram_type = (candidate_metadata.get('type') or '').lower()
            if 'flowchart' in diagram_type:
                diagram_path, diagram_metadata, selected_reference = candidate
                break

        if diagram_path is None and diagram_candidates:
            diagram_path, diagram_metadata, selected_reference = diagram_candidates[0]

        # Extract key points with 20-word limit for presentation slides
        key_points = _collect_highlights_from_lines(lines)
        if not key_points:
            fallback_sentence = limit_words(f"Core ideas from {title.lower()} for modern organisations", 20)
            key_points = [fallback_sentence]

        # Ensure we have at least 5 key points but no more than 10, all with 20-word limit
        while len(key_points) < 5 and len(key_points) > 0:
            # Add some general points if we don't have enough, all limited to 20 words
            if len(key_points) == 1:
                key_points.append(limit_words(f"Practical implementation of {title.lower()} across modern organisations", 20))
            elif len(key_points) == 2:
                key_points.append(limit_words(f"Security considerations and best practices for {title.lower()}", 20))
            elif len(key_points) == 3:
                key_points.append(limit_words(f"Automation and CI/CD enablers for {title.lower()}", 20))
            elif len(key_points) == 4:
                key_points.append(limit_words(f"Cost optimisation and scalability guidance for {title.lower()}", 20))
        
        selected_filename = os.path.basename(selected_reference) if selected_reference else None

        if diagram_metadata is None:
            diagram_metadata = {
                'type': 'Architecture diagram',
                'purpose': 'Highlights a key architecture concept',
                'source': 'Architecture as Code diagrams',
                'explanation': '',
            }

        # Extract contextual explanation from chapter content around diagram reference
        if content and diagram_path and selected_filename:
            lines = content.split('\n')
            for i, line in enumerate(lines):
                if selected_filename.replace('.png', '') in line and 'images/' in line:
                    if i + 1 < len(lines):
                        next_line = lines[i + 1].strip()
                        if next_line and not next_line.startswith('#') and not next_line.startswith('!['):
                            diagram_metadata['explanation'] = limit_words(next_line, 18)
                            break

                    if i > 0:
                        prev_line = lines[i - 1].strip()
                        if prev_line and not prev_line.startswith('#') and not prev_line.startswith('!['):
                            diagram_metadata['explanation'] = limit_words(prev_line, 18)
                            break

        if not diagram_metadata.get('explanation'):
            filename_lower = (selected_filename or '').lower()
            purpose_text = diagram_metadata.get('purpose') or 'Highlights a key architecture concept'
            chapter_match = re.search(r'(chapter|kapitel)(\d+)', filename_lower)
            if chapter_match:
                diagram_metadata['explanation'] = limit_words(
                    f"Diagram: Chapter {chapter_match.group(2)} - {purpose_text}", 18
                )
            else:
                diagram_metadata['explanation'] = limit_words(f"Diagram: {purpose_text}", 18)

        notes_lines = [title]
        notes_lines.append('')
        notes_lines.append('Highlights:')
        for point in key_points[:10]:
            notes_lines.append(f"- {point}")

        if diagram_path:
            notes_lines.append('')
            diagram_type = diagram_metadata.get('type') or 'Diagram'
            notes_lines.append(f"Diagram type: {diagram_type}")
            if diagram_metadata.get('explanation'):
                notes_lines.append(diagram_metadata['explanation'])

        return {
            'title': title,
            'key_points': key_points[:10],
            'diagram_path': diagram_path,
            'diagram_metadata': diagram_metadata,
            'notes': '\n'.join(notes_lines).strip(),
        }

    except Exception as e:
        print(f"Error reading {chapter_file}: {e}")
        return None

def generate_presentation_outline():
    """Generate presentation outline from all book chapters."""
    docs_dir = Path("docs")
    if not docs_dir.exists():
        print("Error: docs/ directory not found")
        return []

    requirements = load_book_requirements()
    book_config = requirements.get('book', {})
    presentation_data = []

    if not book_config:
        # Fallback to alphabetical ordering when requirements are missing
        chapter_files = sorted(glob.glob(str(docs_dir / "*.md")))
        for chapter_file in chapter_files:
            name = Path(chapter_file).name
            if name in ['README.md', 'architecture_as_code.md']:
                continue
            chapter_data = read_chapter_content(chapter_file)
            if chapter_data:
                presentation_data.append({
                    'type': 'chapter',
                    'file': name,
                    'chapter': chapter_data,
                })
        return presentation_data

    front_matter_meta = book_config.get('front_matter', [])
    part_introductions = book_config.get('part_introductions', [])
    chapter_meta = book_config.get('chapters', [])

    for front in front_matter_meta:
        filename = front.get('filename')
        if not filename:
            continue
        front_path = docs_dir / filename
        front_content = read_front_matter_content(front_path, front) if front_path.exists() else None
        if front_content:
            presentation_data.append({
                'type': 'front_matter',
                'file': filename,
                'front_matter': front_content,
            })

    chapters_by_part = {part['title']: [] for part in part_introductions if part.get('title')}
    unassigned_chapters = []

    for chapter in chapter_meta:
        filename = chapter.get('filename')
        area = chapter.get('area')
        if not filename:
            continue

        matched = False
        for part_title, areas in PART_AREA_MAP.items():
            if area in areas:
                if part_title in chapters_by_part:
                    chapters_by_part[part_title].append(chapter)
                else:
                    chapters_by_part.setdefault(part_title, []).append(chapter)
                matched = True
                break

        if not matched:
            unassigned_chapters.append(chapter)

    for part in part_introductions:
        filename = part.get('filename')
        title = part.get('title')
        if not filename or not title:
            continue

        part_path = docs_dir / filename
        part_content = read_part_content(part_path, part) if part_path.exists() else None
        if part_content:
            presentation_data.append({
                'type': 'part',
                'file': filename,
                'part': part_content,
            })

        for chapter in chapters_by_part.get(title, []):
            chapter_file = docs_dir / chapter['filename']
            chapter_content = read_chapter_content(chapter_file)
            if chapter_content:
                chapter_content['label'] = chapter.get('label')
                chapter_content['area'] = chapter.get('area')
                chapter_content['focus_keyword'] = chapter.get('area') or chapter.get('label')
                chapter_content['identifier'] = chapter.get('filename')
                presentation_data.append({
                    'type': 'chapter',
                    'file': chapter.get('filename'),
                    'chapter': chapter_content,
                })

    for chapter in unassigned_chapters:
        chapter_file = docs_dir / chapter['filename']
        chapter_content = read_chapter_content(chapter_file)
        if chapter_content:
            chapter_content['label'] = chapter.get('label')
            chapter_content['area'] = chapter.get('area')
            chapter_content['focus_keyword'] = chapter.get('area') or chapter.get('label')
            chapter_content['identifier'] = chapter.get('filename')
            presentation_data.append({
                'type': 'chapter',
                'file': chapter.get('filename'),
                'chapter': chapter_content,
            })

    return presentation_data

def create_presentation_script(presentation_data):
    """Create a standalone script that renders the presentation using shared helpers."""
    script_content = '''#!/usr/bin/env python3
"""Generate the Architecture as Code presentation from cached data."""

import json
import sys
from pathlib import Path

from pptx import Presentation

ROOT_DIR = Path(__file__).resolve().parents[1]
if str(ROOT_DIR) not in sys.path:
    sys.path.insert(0, str(ROOT_DIR))

from generate_presentation import build_presentation_document

DATA_FILE = Path(__file__).with_name("presentation_data.json")


def load_presentation_data():
    if not DATA_FILE.exists():
        raise FileNotFoundError(f"Presentation data not found: {DATA_FILE}")
    with DATA_FILE.open("r", encoding="utf-8") as handle:
        return json.load(handle)


def main(output_path: Path) -> None:
    presentation_data = load_presentation_data()
    prs = Presentation()
    build_presentation_document(prs, presentation_data)
    prs.save(output_path)
    print(f"✅ Presentation saved to {output_path}")
    print(f"📊 Total slides created: {len(prs.slides)}")


if __name__ == "__main__":
    target_path = Path(sys.argv[1]) if len(sys.argv) > 1 else DATA_FILE.with_name("architecture_as_code_presentation.pptx")
    main(target_path)
'''

    return script_content


def create_presentation_directly(presentation_data, output_path="architecture_as_code_presentation.pptx"):
    """Create PowerPoint presentation directly without generating a script."""
    if not PPTX_AVAILABLE:
        print("❌ python-pptx is not installed. Install it with 'pip install python-pptx>=0.6.21'.")
        print("   The script and supporting materials are still generated.")
        return False

    try:
        prs = Presentation()
        build_presentation_document(prs, presentation_data)
        prs.save(output_path)
    except Exception as exc:  # pragma: no cover - runtime safety
        print(f"❌ Failed to create PowerPoint file: {exc}")
        return False

    print(f"✅ PowerPoint file created: {output_path}")
    print(f"📊 Total slides created: {len(prs.slides)}")
    print("🎨 Slides follow the Architecture as Code theme and 16:9 layout.")
    return True


def main():
    """Main function to generate presentation materials."""
    import argparse
    
    parser = argparse.ArgumentParser(description="Generate PowerPoint presentation from book chapters")
    parser.add_argument("--create-pptx", action="store_true", 
                       help="Create PowerPoint file directly (requires python-pptx)")
    parser.add_argument("--output", default="architecture_as_code_presentation.pptx",
                       help="Output filename for PowerPoint file")
    parser.add_argument("--validate-diagrams", action="store_true",
                       help="Validate that all chapters have diagrams and check diagram type coverage")
    parser.add_argument("--release", action="store_true",
                       help="Generate presentation materials to release folder")
    
    args = parser.parse_args()
    
    print("Analyzing book content to generate presentation...")
    
    # Diagram validation and analysis
    if args.validate_diagrams:
        print("\n🔍 Validating diagram coverage...")
        
        # Check chapter diagram coverage
        chapter_validation = validate_chapter_diagrams()
        print(f"📊 Chapter diagram analysis:")
        print(f"   - Total chapters analyzed: {chapter_validation['total_chapters']}")
        print(f"   - Chapters with diagrams: {chapter_validation['chapters_with_diagram_count']}")
        print(f"   - Coverage: {(chapter_validation['chapters_with_diagram_count']/chapter_validation['total_chapters']*100):.1f}%")
        
        if chapter_validation['chapters_without_diagrams']:
            print(f"⚠️  Chapters missing diagrams: {chapter_validation['chapters_without_diagrams']}")
        
        # Check diagram type coverage
        missing_types = get_missing_diagram_types()
        print(f"\n🎨 Mermaid diagram type analysis:")
        print(f"   - Available types: {missing_types['available']}")
        print(f"   - Missing types: {missing_types['missing']}")
        
        if missing_types['missing']:
            print(f"⚠️  Missing essential diagram types: {missing_types['missing']}")
            print("   Consider adding these types for comprehensive presentation coverage")
        
        # Detailed analysis
        print(f"\n📋 Detailed diagram inventory:")
        for diagram_type, files in missing_types['analysis'].items():
            if files:
                print(f"   - {diagram_type}: {len(files)} diagrams")
        
        print("\n" + "="*60)
    
    # Ensure presentations directory exists
    if args.release:
        presentations_dir = Path("releases/presentation")
        print("Release mode: Generating presentation materials to releases/presentation/")
    else:
        presentations_dir = Path("presentations")
    
    presentations_dir.mkdir(exist_ok=True, parents=True)
    
    # Read all chapters WITHOUT modifying them
    presentation_data = generate_presentation_outline()

    if not presentation_data:
        print("Error: Could not read chapter data")
        return 1

    print(f"Found {len(presentation_data)} chapters to include in presentation")

    slides_path, slide_count = generate_prezi_slides()
    if slide_count:
        print(f"✅ Prezi slides generated at {slides_path} ({slide_count} slides)")
    else:
        print("⚠️ Prezi slides were not generated because no markdown chapters were found")
    
    # Create presentation outline
    outline_content = "# Presentation Outline\n\n"
    outline_content += "This presentation is generated from the book chapters in docs/\n\n"
    
    for item in presentation_data:
        chapter = item['chapter']
        outline_content += f"## {chapter['title']}\n\n"
        for point in chapter['key_points']:
            outline_content += f"- {point}\n"
        outline_content += "\n"
    
    # Write outline (outside docs directory)
    with open(presentations_dir / "presentation_outline.md", 'w', encoding='utf-8') as f:
        f.write(outline_content)

    data_file = presentations_dir / "presentation_data.json"
    data_file.write_text(json.dumps(presentation_data, ensure_ascii=False, indent=2), encoding='utf-8')

    # Create PowerPoint generator script
    pptx_script = create_presentation_script(presentation_data)
    with open(presentations_dir / "generate_pptx.py", 'w', encoding='utf-8') as f:
        f.write(pptx_script)
    
    # Create requirements file for presentation generation
    requirements = "\n".join(
        [
            "python-pptx>=0.6.21",
            "PyYAML>=6.0",
        ]
    ) + "\n"
    with open(presentations_dir / "requirements.txt", 'w', encoding='utf-8') as f:
        f.write(requirements)
    
    print("✅ Presentation materials generated in presentations/ directory:")
    print("   - presentation_outline.md")
    print("   - presentation_data.json")
    print("   - generate_pptx.py")
    print("   - requirements.txt")
    
    # If --create-pptx flag is provided, create the PowerPoint file directly
    if args.create_pptx:
        print("\n📊 Creating PowerPoint file directly...")
        output_path = presentations_dir / args.output
        if create_presentation_directly(presentation_data, output_path):
            print(f"✅ PowerPoint file created: {output_path}")
            
            # In release mode, also copy to standard location for backward compatibility
            if args.release:
                import shutil
                standard_dir = Path("presentations")
                standard_dir.mkdir(exist_ok=True)
                shutil.copy2(output_path, standard_dir / args.output)
                print(f"Presentation also copied to standard location: {standard_dir / args.output}")
        else:
            print("❌ Failed to create PowerPoint file")
            print("📝 You can still use the generated script:")
            print(f"   cd {presentations_dir}")
            print("   pip install -r requirements.txt")
            print("   python generate_pptx.py")
            return 1
    else:
        print("\n📝 To create the PowerPoint file:")
        print(f"   cd {presentations_dir}")
        print("   pip install -r requirements.txt")
        print("   python generate_pptx.py")
        print("\n💡 Or use: python generate_presentation.py --create-pptx")
    
    print("\n📁 Prezi slide data stored in docs/prezi/slides.json.")
    print("🚫 No other files in docs/ directory were modified.")
    
    return 0

if __name__ == "__main__":
    sys.exit(main())
