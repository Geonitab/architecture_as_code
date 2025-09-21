#!/usr/bin/env python3
"""
PowerPoint Presentation Generator for Arkitektur som kod
Generated automatically from book content.
Complies with Swedish standards and book theme.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

def create_presentation():
    """Create PowerPoint presentation with book content."""
    prs = Presentation()
    
    # Set up Swedish theme colors (inspired by Swedish flag and professional standards)
    # Blue: #006AA7 (Swedish blue), Yellow: #FECC00 (Swedish yellow), Gray: #333333
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Arkitektur som kod"
    subtitle.text = "En omfattande guide för svenska organisationer"
    
    # Style the title slide
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
    

    # Chapter: Inledning till arkitektur som kod
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Inledning till arkitektur som kod"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_01_inledning.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Från Infrastructure as Code till Architecture as Code: Infrastructure as Code (IaC) var det första steget mot kodifiering av IT-resurser. Genom att behandla infrastruktur som kod uppnåddes automatiser..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Definition och omfattning: Architecture as Code definieras som praktiken att beskriva, versionhantera och automatisera hela systemarkitekturen genom maskinläsbar kod. Detta omfattar inte bara infrastr..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Bokens syfte och målgrupp: Denna bok vänder sig till systemarkitekter, utvecklare, devops-ingenjörer och projektledare som vill förstå och implementera Architecture as Code i sina organisationer. Infr..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Säkerhetsaspekter inom inledning till arkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automatisering och CI/CD för inledning till arkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadsoptimering inom inledning till arkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för inledning till arkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av inledning till arkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för inledning till arkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för inledning till arkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Grundläggande principer för Architecture as Code
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Grundläggande principer för Architecture as Code"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_02_kapitel1.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Deklarativ arkitekturdefinition: Den deklarativa approachen inom Architecture as Code innebär att beskriva önskat systemtillstånd på alla nivåer - från applikationskomponenter till infrastruktur. Dett..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Helhetsperspektiv på kodifiering: Medan Infrastructure as Code fokuserar på infrastrukturresurser, omfattar Architecture as Code hela systemekosystemet. Detta inkluderar applikationslogik, dataflöden,..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Immutable architecture patterns: Principen om immutable arkitektur bygger vidare på Infrastructure as Code:s immutable infrastruktur men applicerar det på hela systemarkitekturen. Istället för att mod..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Testbarhet på arkitekturnivå: Architecture as Code möjliggör testning av hela systemarkitekturen, inte bara enskilda komponenter. Detta inkluderar validering av arkitekturmönster, compliance med desig..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automatisering och CI/CD för grundläggande principer för architecture as code"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadsoptimering inom grundläggande principer för architecture as code"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för grundläggande principer för architecture as code"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av grundläggande principer för architecture as code"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för grundläggande principer för architecture as code"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för grundläggande principer för architecture as code"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Versionhantering och kodstruktur
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Versionhantering och kodstruktur"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_03_kapitel2.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Git-baserad arbetsflöde för infrastruktur: Git utgör standarden för versionhantering av IaC-kod och möjliggör distribuerat samarbete mellan team-medlemmar. Varje förändring dokumenteras med commit-med..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kodorganisation och modulstruktur: Välorganiserad kodstruktur är avgörande för maintainability och collaboration i större IaC-projekt. Modulär design möjliggör återanvändning av infrastrukturkomponent..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Svenska organisationers perspektiv på versionhantering och kodstruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Säkerhetsaspekter inom versionhantering och kodstruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automatisering och CI/CD för versionhantering och kodstruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadsoptimering inom versionhantering och kodstruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för versionhantering och kodstruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av versionhantering och kodstruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för versionhantering och kodstruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för versionhantering och kodstruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Architecture Decision Records (ADR)
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Architecture Decision Records (ADR)"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_04_adr_kapitel.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Övergripande beskrivning: Architecture Decision Records (ADR) utgör en systematisk approach för att dokumentera viktiga arkitekturbeslut som påverkar systemets struktur, prestanda, säkerhet och underh..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Vad är Architecture Decision Records?: Architecture Decision Records definieras som korta textdokument som fångar viktiga arkitekturbeslut tillsammans med deras kontext och konsekvenser. Varje ADR bes..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Struktur och komponenter av ADR: Varje ADR följer en konsekvent struktur som säkerställer att all relevant information fångas systematiskt: ADR numreras sekventiellt (ADR-0001, ADR-0002, etc.) för att..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Verktyg och best practices för ADR: Flera verktyg underlättar creation och management av ADR: **adr-tools**: Command-line verktyg för att skapa och hantera ADR-filer **adr-log**: Automatisk generering..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Integration med Architecture as Code: ADR spelar en central roll i Architecture as Code-metodik genom att dokumentera designbeslut som sedan implementeras som kod. Denna integration skapar en tydlig k..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Compliance och kvalitetsstandarder: ADR-metodik stödjer svenska compliance-krav genom strukturerad dokumentation som möjliggör: **Regulatory Compliance**: Systematisk dokumentation för GDPR, PCI-DSS o..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtida utveckling och trends: ADR-metodik utvecklas kontinuerligt med integration av nya verktyg och processer: **AI-assisterade ADR**: Machine learning för att identifiera när nya ADR behövs basera..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Architecture Decision Records representerar en fundamental komponent i modern Architecture as Code-metodik. Genom strukturerad dokumentation av arkitekturbeslut skapas transparens, spå..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för architecture decision records (adr)"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för architecture decision records (adr)"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Automatisering och CI/CD-pipelines
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Automatisering och CI/CD-pipelines"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_04_kapitel3.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Den teoretiska grunden för CI/CD-automation: Continuous Integration och Continuous Deployment representerar mer än bara tekniska processer - de utgör en filosofi för mjukvaruutveckling som prioriterar..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• CI/CD-fundamentals för svenska organisationer: Svenska organisationer opererar i en komplex regulatorisk miljö som kräver särskild uppmärksamhet vid implementering av CI/CD-pipelines för Infrastructur..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Pipeline design principles: Effektiva CI/CD-pipelines för Infrastructure as Code bygger på fundamentala design principles som optimerar för speed, safety och observability. Dessa principles måste anpa..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadsoptimering och budgetkontroll: Svenska organisationer måste hantera infrastrukturkostnader med particular attention till valutafluktuationer, regional pricing variations och compliance-relater..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Pipeline design principles: Effektiva IaC-pipelines bygger på principerna för fail-fast feedback och progressive deployment. Tidiga valideringssteg identifierar problem innan kostsamma infrastrukturfö..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automated testing strategier: Multi-level testing strategies för IaC inkluderar syntax validation, unit testing av moduler, integration testing av komponenter, och end-to-end testing av kompletta milj..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Infrastructure validation: Pre-deployment validation säkerställer att infrastrukturändringar möter organisatoriska requirements innan de appliceras. Detta inkluderar policy compliance, security postur..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Deployment strategier: Blue-green deployments och canary releases anpassas för infrastrukturkontext genom att skapa parallella miljöer eller successivt rulla ut förändringar. Rolling deployments hante..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Monitoring och observability: Pipeline observability inkluderar både execution metrics och business impact measurements. Technical metrics som build time, success rate, och deployment frequency kombin..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Automatisering och CI/CD-pipelines för Infrastructure as Code utgör en kritisk komponent för svenska organisationer som strävar efter digital excellence och regulatory compliance. Geno..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: DevOps och CI/CD för Architecture as Code
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "DevOps och CI/CD för Architecture as Code"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_08_kapitel6.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: DevOps och CI/CD för Architecture as Code skapar grunden för helhetlig, kodifierad systemutveckling där alla aspekter av arkitekturen hanteras som kod. För svenska organisationer inneb..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Från Infrastructure as Code till Architecture as Code DevOps: Traditionella DevOps-praktiker fokuserade primärt på applikationsutveckling och deployment, medan Infrastructure as Code (IaC) utvidgade d..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• CI/CD-pipelines för Architecture as Code: Architecture as Code CI/CD-pipelines skiljer sig från traditionella pipelines genom att hantera flera sammankopplade arkitekturdomäner samtidigt. Istället för..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Architecture as Code Testing-strategier: Architecture as Code kräver testing-strategier som går beyond traditionell infrastruktur- eller applikationstestning. Testning måste validera arkitekturkonsist..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Architecture as Code Monitoring och Observability: Architecture as Code monitoring måste hantera komplexiteten av att övervaka hela arkitekturlandskapet som en sammanhängande enhet. Detta inkluderar t..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• DevOps Kultur för Architecture as Code: Architecture as Code kräver en mogen DevOps-kultur som kan hantera komplexiteten av holistic systemtänkande samtidigt som den bibehåller agilitet och innovation..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Referenser och vidare läsning: - Datainspektionen. \"GDPR för svenska organisationer.\" Vägledning om personuppgiftsbehandling. - Myndigheten för samhällsskydd och beredskap (MSB). \"Säkerhetsskydd fö..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av devops och ci/cd för architecture as code"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för devops och ci/cd för architecture as code"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för devops och ci/cd för architecture as code"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Molnarkitektur som kod
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Molnarkitektur som kod"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_05_kapitel4.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Molnleverantörers ekosystem för IaC: Svenska organisationer står inför ett rikt utbud av molnleverantörer, var och en med sina egna styrkor och specialiseringar. För att uppnå framgångsrik cloud adopt..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Cloud-native IaC patterns: Cloud-native Infrastructure as Code patterns utnyttjar molnspecifika tjänster och capabilities för att skapa optimala arkitekturer. Dessa patterns inkluderar serverless comp..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Multi-cloud strategier: Multi-cloud Infrastructure as Code strategier möjliggör distribution av workloads across flera molnleverantörer för att optimera kostnad, prestanda, och resiliens. Provider-agn..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Serverless infrastruktur: Serverless Infrastructure as Code fokuserar på function definitions, event triggers, och managed service configurations istället för traditionell server management. Detta app..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Praktiska implementationsexempel: För att demonstrera molnarkitektur som kod i praktiken för svenska organisationer, presenteras här kompletta implementationsexempel som visar how real-world scenarios..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Molnarkitektur som kod representerar en fundamental evolution av Infrastructure as Code för svenska organisationer som opererar i cloud-native miljöer. Genom att utnyttja cloud provide..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för molnarkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av molnarkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för molnarkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för molnarkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Containerisering och orkestrering som kod
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Containerisering och orkestrering som kod"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_11_kapitel10.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Container-teknologiens roll inom IaC: Containers erbjuder application-level virtualization som paketerar applikationer med alla dependencies i isolated, portable units. För Infrastructure as Code inne..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kubernetes som orchestration platform: Kubernetes har emergerat som leading container orchestration platform genom dess declarative configuration model och extensive ecosystem. YAML-based manifests de..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Service mesh och advanced networking: Service mesh architectures som Istio och Linkerd implementeras through Infrastructure as Code för att hantera inter-service communication, security policies, och ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Infrastructure automation med container platforms: Container-native infrastructure tools som Crossplane och Operator Framework extend Kubernetes för complete infrastructure management. These platforms..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Persistent storage och data management: Persistent volume management för containerized applications kräver careful consideration av performance, availability, och backup requirements. Storage classes ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Containerisering och orkestrering som kod transformerar application deployment från manual, error-prone processes till automated, reliable workflows. Kubernetes och associerade verktyg..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - Kubernetes Documentation. \"Concepts and Architecture.\" The Kubernetes Project. - Docker Inc. \"Docker Best Practices.\" Docker Documentation. - Cloud Native Computing Founda..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av containerisering och orkestrering som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för containerisering och orkestrering som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för containerisering och orkestrering som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Microservices-arkitektur som kod
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Microservices-arkitektur som kod"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_13_kapitel12.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Den evolutionära resan från monolit till microservices: Svenska företag som Spotify, Klarna, King och H&M har blivit globala digitala ledare genom att anta microservices-arkitektur tidigt. Deras framg..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Microservices design principles för IaC: Att framgångsrikt implementera microservices-arkitektur kräver en djup förståelse för de designprinciper som styr både service-design och infrastrukturen som s..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Service discovery och communication patterns: I en microservices-arkitektur är förmågan för tjänster att hitta och kommunicera med varandra fundamental för systemets funktionalitet. Service discovery ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Data management i distribuerade system: En av de mest fundamentala utmaningarna i microservices-arkitektur är hur data ska hanteras och delas mellan tjänster. Traditional monolithic applications har t..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Service mesh implementation: Service mesh technology representerar en paradigm shift i hur microservices kommunicerar och hanterar cross-cutting concerns. Istället för att implementera communication l..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Deployment och scaling strategies: Modern microservices-arkitektur kräver sophisticated deployment och scaling strategies som kan hantera hundreds eller thousands av independent services. För svenska ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Monitoring och observability: I en microservices-arkitektur där requests kan traverse dozens av services blir traditional monitoring approaches inadequate. Comprehensive observability blir essential f..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Microservices-arkitektur som kod representerar mer än bara en teknisk evolution - det är en transformation som påverkar hela organisationen, från hur team organiseras till hur affärspr..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - Martin Fowler. \"Microservices Architecture.\" Martin Fowler\'s Blog. - Netflix Technology Blog. \"Microservices at Netflix Scale.\" Netflix Engineering. - Kubernetes Document..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för microservices-arkitektur som kod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Säkerhet i Architecture as Code
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Säkerhet i Architecture as Code"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_06_kapitel5.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Säkerhetsarkitekturens dimensioner: *Mindmappen illustrerar de komplexa sambanden mellan olika säkerhetsaspekter i Architecture as Code, från threat modeling och Zero Trust Architecture till Policy as..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kapitelets omfattning och mål: Säkerhetsutmaningarna i dagens digitala landskap kräver en fundamental omvärdering av traditionella säkerhetsapproacher. När organisationer adopterar Architecture as Cod..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Teoretisk grund: Säkerhetsarkitektur i den digitala tidsåldern: Den traditionella säkerhetsfilosofin byggde på förutsättningen om en tydlig gräns mellan \"insidan\" och \"utsidan\" av organisationen. ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Policy as Code: Automatiserad säkerhetsstyrning: Traditionell säkerhetsstyrning bygger på manuella processer, dokumentbaserade policies och människodrivna kontroller. Säkerhetsavdelningar författar po..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Security-by-design: Arkitektoniska säkerhetsprinciper: Security-by-design representerar inte bara en implementationsstrategi utan en fundamental filosofisk approach till systemarkitektur. Traditionell..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Policy as Code implementation: Policy as Code representerar paradigmskiftet från manuella säkerhetspolicies till automatiserat policy enforcement genom programmatiska definitioner. Open Policy Agent (..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Secrets Management och Data Protection: Modern distributed architectures proliferate secrets exponentially compared till traditional monolithic applications. API keys, database credentials, encryption..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Secrets management och data protection: Comprehensive secrets management utgör foundationen för säker IaC implementation. Secrets som API keys, databas-credentials och encryption keys måste hanteras g..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Nätverkssäkerhet och microsegmentering: Traditional network security architectures built på assumption av trusted internal networks separated från untrusted external networks through perimeter defense..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Avancerade Säkerhetsarkitekturmönster: Modern enterprise säkerhetsarkitekturer kräver sofistikerad orchestration av multiple security tools och processes för att hantera växande volymer av security ev..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Policy och säkerhet som kod i detalj
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Policy och säkerhet som kod i detalj"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_12_kapitel11.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Introduktion och kontextualisering: I en värld där svenska organisationer hanterar allt mer komplexa digitala infrastrukturer samtidigt som regulatoriska krav skärps kontinuerligt, har Policy as Code ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Evolutionen av säkerhetshantering inom Infrastructure as Code: Säkerhetshantering inom Infrastructure as Code har genomgått en betydande evolution från ad-hoc skript och manuella checklistor till sofi..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Open Policy Agent (OPA) och Rego: Grunden för policy-driven säkerhet: Open Policy Agent har etablerats som de facto standarden för policy as code implementation genom sin flexibla arkitektur och kraft..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• OSCAL: Open Security Controls Assessment Language - Revolutionerande säkerhetsstandardisering: Open Security Controls Assessment Language (OSCAL) representerar en paradigmskifte inom säkerhetshanterin..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Gatekeeper och Kubernetes Policy Enforcement: Enterprise-grade implementationer: Kubernetes-miljöer representerar en unik utmaning för policy enforcement på grund av deras dynamiska natur och complex ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadsoptimering inom policy och säkerhet som kod i detalj"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för policy och säkerhet som kod i detalj"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av policy och säkerhet som kod i detalj"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för policy och säkerhet som kod i detalj"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för policy och säkerhet som kod i detalj"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Compliance och regelefterlevnad
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Compliance och regelefterlevnad"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_12_compliance.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• AI och maskininlärning för infrastrukturautomatisering: Artificiell intelligens revolutionerar Infrastructure as Code genom intelligent automation, prediktiv skalning och självläkande system. Maskinin..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Cloud-native och serverless utveckling: Serverless computing fortsätter att utvecklas bortom enkla function-as-a-service mot omfattande serverless-arkitekturer. Infrastructure as Code måste anpassas f..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Policydriven infrastruktur och styrning: Policy as Code blir allt mer sofistikerat med automatiserad compliance-kontroll, kontinuerlig styrningsverkställighet och dynamisk policyanpassning. Policyer u..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kvantdatorer och nästa generations teknologier: Kvantdatorers påverkan på Infrastructure as Code kommer att kräva en grundläggande omtänkning av säkerhetsmodeller, beräkningsarkitekturer och resurshan..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Hållbarhet och grön databehandling: Miljöhållbarhet blir central övervägande för infrastrukturdesign och drift. Kolmedveten infrastrukturhantering skiftar automatiskt arbetsbelastningar till regioner ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Framtida Infrastructure as Code-utveckling kommer att drivas av AI-automation, serverless-arkitekturer, beredskap för kvantdatorer och hållbarhetskrav. Organisationer måste proaktivt i..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - IEEE Computer Society. \"Quantum Computing Impact on Infrastructure.\" IEEE Quantum Computing Standards. - Green Software Foundation. \"Sustainable Infrastructure Patterns.\" ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av compliance och regelefterlevnad"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för compliance och regelefterlevnad"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för compliance och regelefterlevnad"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Teststrategier för infrastruktukod
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Teststrategier för infrastruktukod"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_17_kapitel16.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Övergripande beskrivning: Testning av Infrastructure as Code skiljer sig fundamentalt från traditionell mjukvarutestning genom att fokusera på infrastrukturkonfiguration, resurskompatibilitet och milj..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Unit testing för infrastrukturkod: Unit testing för Infrastructure as Code fokuserar på validation av enskilda moduler och resources utan att faktiskt deploya infrastruktur. Detta möjliggör snabb feed..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Integrationstesting och miljövalidering: Integration testing för Infrastructure as Code verifierar att different infrastructure components fungerar tillsammans korrekt och att deployed infrastruktur m..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Security och compliance testing: Security testing för Infrastructure as Code måste validate både infrastructure configuration security och operational security controls. Detta inkluderar scanning för ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Performance och skalbarhetstesting: Performance testing för Infrastructure as Code fokuserar på validation av infrastructure capacity, response times och resource utilization under various load condit..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Krav som kod och testbarhet: *Relationen mellan affärskrav, funktionella krav och verifieringsmetoder illustrerar hur Infrastructure as Code möjliggör spårbar testning från högre abstraktionsnivåer ne..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Comprehensive testing strategies för Infrastructure as Code är essential för att säkerställa reliable, secure och cost-effective infrastructure deployments. En väl designad test pyrami..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - Terratest Documentation. \"Infrastructure Testing for Terraform.\" Gruntwork, 2023. - Open Policy Agent. \"Policy Testing Best Practices.\" CNCF OPA Project, 2023. - AWS. \"In..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för teststrategier för infrastruktukod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för teststrategier för infrastruktukod"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Architecture as Code i praktiken
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Architecture as Code i praktiken"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_08_kapitel7.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Implementation roadmap och strategier: Successful Architecture as Code adoption följer vanligen en phased approach som börjar med pilot projects och gradvis expanderar till enterprise-wide implementat..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Tool selection och ecosystem integration: Technology stack selection balanserar organizational requirements med market maturity och community support. Terraform har emerged som leading multi-cloud sol..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Production readiness och operational excellence: Security-first approach implementerar comprehensive security controls från design phase. Secrets management, access controls, audit logging, och compli..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Common challenges och troubleshooting: State management complexity grows significantly som infrastructure scales och involves multiple teams. State file corruption, concurrent modifications, och state..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Enterprise integration patterns: Multi-account/subscription strategies för cloud environments provide isolation, security boundaries, och cost allocation capabilities. Infrastructure code måste handle..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Practical Infrastructure as Code implementation balanserar technical excellence med organizational realities. Success kräver comprehensive planning, stakeholder alignment, incremental ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - HashiCorp. \"Terraform Best Practices.\" HashiCorp Learn Platform. - AWS Well-Architected Framework. \"Infrastructure as Code.\" Amazon Web Services. - Google Cloud. \"Infrast..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av architecture as code i praktiken"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för architecture as code i praktiken"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för architecture as code i praktiken"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Kostnadsoptimering och resurshantering
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Kostnadsoptimering och resurshantering"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_16_kapitel15.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Övergripande beskrivning: Kostnadsoptimering utgör en kritisk komponent i Infrastructure as Code-implementationer, särskilt när organisationer migrerar till molnbaserade lösningar. Utan proper cost ma..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• FinOps och cost governance: FinOps representerar en växande disciplin som kombinerar finansiell hantering med molnoperationer för att maximera affärsvärdet av molninvesteringar. Inom IaC-kontext inneb..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automatisk resursskalning och rightsizing: Automatisk resursskalning utgör kärnan i kostnadseffektiv Infrastructure as Code. Genom att definiera skalningsregler baserade på faktiska användningsmönster..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Cost monitoring och alerting: Comprehensive cost monitoring kräver integration av monitoring-verktyg direkt i IaC-konfigurationerna. CloudWatch, Azure Monitor och Google Cloud Monitoring kan konfigure..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Multi-cloud cost optimization: Multi-cloud strategier kompliserar kostnadsoptimering men erbjuder också möjligheter för cost arbitrage mellan olika leverantörer. IaC-verktyg som Terraform möjliggör co..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Kostnadsoptimering inom Infrastructure as Code kräver systematisk approach som kombinerar tekniska verktyg, automatiserade processer och organisatorisk medvetenhet. Framgångsrik implem..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - AWS. \"AWS Cost Optimization Guide.\" Amazon Web Services Documentation, 2023. - FinOps Foundation. \"FinOps Framework och Best Practices.\" The Linux Foundation, 2023. - Kube..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av kostnadsoptimering och resurshantering"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för kostnadsoptimering och resurshantering"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för kostnadsoptimering och resurshantering"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Migration från traditionell infrastruktur
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Migration från traditionell infrastruktur"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_18_kapitel17.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Övergripande beskrivning: Migration från traditionell, manuellt konfigurerad infrastruktur till Infrastructure as Code representerar en av de mest kritiska transformationerna för moderna IT-organisati..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Assessment och planning faser: Comprehensive infrastructure assessment utgör foundationen för successful IaC migration. Detta inkluderar inventory av existing resources, dependency mapping, risk asses..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Lift-and-shift vs re-architecting: Lift-and-shift migration representerar den snabbaste vägen till cloud adoption men limiterar potential benefits från cloud-native capabilities. Denna approach är läm..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Gradvis kodifiering av infrastruktur: Infrastructure inventory automation genom tools som Terraform import, CloudFormation drift detection och Azure Resource Manager templates enables systematic conve..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team transition och kompetensutveckling: Skills development programs måste prepare traditional system administrators och network engineers för IaC-based workflows. Training curricula ska encompass Inf..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Migration från traditionell infrastruktur till Infrastructure as Code representerar en kritisk transformation som kräver systematisk planering, gradvis implementation och omfattande te..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - AWS. \"Large-Scale Migration och Modernization Guide.\" Amazon Web Services, 2023. - Microsoft. \"Azure Migration Framework och Best Practices.\" Microsoft Azure Documentation..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av migration från traditionell infrastruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för migration från traditionell infrastruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för migration från traditionell infrastruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Organisatorisk förändring och teamstrukturer
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Organisatorisk förändring och teamstrukturer"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_10_kapitel9.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Organisatoriska förändringsprocessens komplexitet: *Mindmappen visualiserar de mångsidiga aspekterna av organisatorisk förändring vid IaC-implementation. Den visar hur DevOps-kulturtransformation, cro..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övergripande beskrivning: Implementering av Infrastructure as Code kräver djupgående organisatoriska förändringar som sträcker sig långt bortom teknisk transformation. Traditionella IT-organisationer ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• DevOps-kulturtransformation: DevOps representerar fundamental kulturförändering från \"us vs them\" mentalitet mellan development och operations till shared ownership av product lifecycle. Denna trans..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Cross-funktionella team strukturer: Cross-functional teams för IaC implementation måste include diverse skills covering software development, systems administration, security engineering och product m..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kompetenshöjning och utbildning: Comprehensive training program för IaC adoption måste cover technical skills, process changes och cultural transformation aspects. Multi-modal learning approaches incl..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Rollförändring och karriärutveckling: Traditional system administrator roles evolve toward Infrastructure Engineers som combine operational expertise med software development skills. Career developmen..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Change management strategier: Change management för IaC adoption måste address both technical och cultural aspects av organizational transformation. Successful change strategies include stakeholder en..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Organisatorisk förändring utgör den mest kritiska komponenten för successful Infrastructure as Code adoption. Technical tools och processes kan implementeras relativt snabbt, men cultu..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - Puppet. \"State of DevOps Report.\" Puppet Labs, 2023. - Google. \"DORA State of DevOps Research.\" Google Cloud, 2023. - Spotify. \"Spotify Engineering Culture.\" Spotify Tec..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för organisatorisk förändring och teamstrukturer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Team-struktur och kompetensutveckling för IaC
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Team-struktur och kompetensutveckling för IaC"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_18_team_struktur.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Organisatorisk transformation för IaC: Traditionella organisationsstrukturer med separata utvecklings-, test- och drift-teams skapar silos som hindrar effektiv Infrastructure as Code adoption. Cross-f..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kompetensområden för IaC-specialister: Infrastructure as Code professionals behöver hybrid skills som kombinerar traditional systems administration knowledge med software engineering practices. Progra..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Utbildningsstrategier och certifieringar: Strukturerade utbildningsprogram kombinerar theoretical learning med hands-on practice för effective skill development. Online platforms som A Cloud Guru, Plu..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Agile team models för infrastructure: Cross-functional infrastructure teams inkluderar cloud engineers, automation specialists, security engineers, och site reliability engineers som collaborerar on s..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kunskapsdelning och communities of practice: Documentation strategies för Infrastructure as Code inkluderar architecture decision records, runbooks, troubleshooting guides, och best practices reposito..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Performance management och career progression: Technical career ladders för Infrastructure as Code specialists provide clear advancement paths from junior automation engineers to senior architect role..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Successful Infrastructure as Code adoption kräver omfattande organisatorisk förändring som går beyond teknisk implementation. Team-strukturer måste redesignas för cross-functional coll..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - Gene Kim, Jez Humble, Patrick Debois, John Willis. \"The DevOps Handbook.\" IT Revolution Press. - Matthew Skelton, Manuel Pais. \"Team Topologies: Organizing Business and Tec..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för team-struktur och kompetensutveckling för iac"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för team-struktur och kompetensutveckling för iac"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Digitalisering genom kodbaserad infrastruktur
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Digitalisering genom kodbaserad infrastruktur"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_19_digitalisering.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Svenska digitaliseringslandskapet: *Mindmappen belyser de unika aspekterna av digitalisering i svensk kontext, från regulatoriska utmaningar och framgångsexempel till de specifika fördelar som IaC erb..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övergripande beskrivning: Digitalisering handlar inte enbart om att införa ny teknik, utan om en fundamental förändring av hur organisationer levererar värde till sina kunder och intressenter. Infrast..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Cloud-first strategier för svensk digitalisering: Sverige har utvecklat en stark position inom molnteknologi, delvis drivet av ambitiösa digitaliseringsmål inom både offentlig och privat sektor samt u..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automatisering av affärsprocesser: IaC möjliggör automatisering som sträcker sig långt bortom traditionell IT-drift till att omfatta hela affärsprocesser med särskild hänsyn till svenska organisatione..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Digital transformation i svenska organisationer: Svenska organisationer genomgår för närvarande en av de mest omfattande digitaliseringsprocesserna i modern tid. Infrastructure as Code utgör ofta den ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Digitalisering genom kodbaserad infrastruktur representerar en fundamental förändring i hur svenska organisationer levererar IT-tjänster och skapar affärsvärde. IaC möjliggör den flexi..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - Digitaliseringsstyrelsen. \"Digitaliseringsstrategi för Sverige.\" Regeringskansliet, 2022. - McKinsey Digital. \"Digital Transformation in the Nordics.\" McKinsey & Company, ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av digitalisering genom kodbaserad infrastruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för digitalisering genom kodbaserad infrastruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för digitalisering genom kodbaserad infrastruktur"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Kapitel 20: Använd Lovable för att skapa mockups för svenska organisationer
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Kapitel 20: Använd Lovable för att skapa mockups för svenska organisationer"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_21_kapitel20.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Inledning till Lovable: Lovable är en AI-driven utvecklingsplattform som revolutionerar hur svenska organisationer kan skapa interaktiva mockups och prototyper. Genom att kombinera naturlig språkbehan..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Steg-för-steg guide för implementering i svenska organisationer: **1. Miljöförberedelse** **2. Svensk lokaliseringskonfiguration** **3. Definiera svenska användarresor**"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning och nästa steg: Lovable erbjuder svenska organisationer en kraftfull plattform för att skapa compliance-medvetna mockups och prototyper. Genom att integrera svenska e-legitimationstjäns..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Säkerhetsaspekter inom kapitel 20: använd lovable för att skapa mockups för svenska organisationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automatisering och CI/CD för kapitel 20: använd lovable för att skapa mockups för svenska organisationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadsoptimering inom kapitel 20: använd lovable för att skapa mockups för svenska organisationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för kapitel 20: använd lovable för att skapa mockups för svenska organisationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av kapitel 20: använd lovable för att skapa mockups för svenska organisationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för kapitel 20: använd lovable för att skapa mockups för svenska organisationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för kapitel 20: använd lovable för att skapa mockups för svenska organisationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Framtida trender och teknologier
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Framtida trender och teknologier"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_19_kapitel18.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Övergripande beskrivning: Infrastructure as Code står inför omfattande transformation driven av teknologiska genombrott inom artificiell intelligens, quantum computing, edge computing och miljömedvete..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Artificiell intelligens och maskininlärning integration: AI och ML-integration i Infrastructure as Code transformerar från reaktiva till prediktiva system som kan anticipera och förebygga problem inna..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Edge computing och distribuerad infrastruktur: Edge computing förändrar fundamentalt hur Infrastructure as Code designas och implementeras. Istället för centraliserade molnresurser distribueras comput..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sustainability och green computing: Environmental sustainability blir allt viktigare inom Infrastructure as Code med fokus på carbon footprint reduction, renewable energy usage och resource efficiency..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Nästa generations IaC-verktyg och paradigm: DevOps evolution fortsätter med nya verktyg och methodologies som förbättrar utvecklarhastighet, operational efficiency och system reliability. GitOps, Plat..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Quantum computing påverkan på säkerhet: Quantum computing development hotar current cryptographic standards och kräver proactive preparation för post-quantum cryptography transition. Infrastructure as..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Sammanfattning: Framtiden för Infrastructure as Code karakteriseras av intelligent automation, environmental sustainability och enhanced security capabilities. Svenska organisationer som investerar i ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Källor och referenser: - NIST. \"Post-Quantum Cryptography Standards.\" National Institute of Standards and Technology, 2024. - IEA. \"Digitalization and Energy Efficiency.\" International Energy Agen..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för framtida trender och teknologier"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för framtida trender och teknologier"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Best practices och lärda läxor
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Best practices och lärda läxor"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_20_kapitel19.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Best practices holistiska perspektiv: *Mindmappen presenterar det omfattande landskapet av best practices och lärda läxor inom Infrastructure as Code. Den visar sambanden mellan kodorganisation, säker..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övergripande beskrivning: Infrastructure as Code best practices representerar culminationen av collective wisdom från tusentals organisationer som har genomgått IaC transformation över det senaste dec..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kod organisation och modulstruktur: Effective code organization utgör foundationen för maintainable och scalable Infrastructure as Code implementations. Well-structured repositories med clear hierarch..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Säkerhet och compliance patterns: Security-first design patterns have emerged as fundamental requirements för modern Infrastructure as Code implementations. These patterns emphasize defense-in-depth, ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Performance och skalning strategier: Infrastructure performance optimization patterns focus på cost efficiency, resource utilization och response time optimization. Swedish e-commerce companies have d..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Governance och policy enforcement: Governance frameworks för Infrastructure as Code must balance developer autonomy med organizational control through clear policies, automated enforcement och excepti..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Internationella erfarenheter och svenska bidrag: Global best practice evolution has been significantly influenced av svenska innovations i organizational design, environmental consciousness och collab..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Incident management och response patterns: Effektiv incidenthantering utgör en kritisk komponent för operational excellence inom Infrastructure as Code-miljöer. När infrastruktur definieras som kod, k..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Dokumentation och knowledge management: Comprehensive documentation strategies för Infrastructure as Code environments must balance technical detail med accessibility för diverse stakeholders. Effecti..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Utbildning och kompetensutveckling: Strategisk kompetensutveckling för Infrastructure as Code requires comprehensive training programs som address both technical skills och organizational transformati..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Slutsats
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Slutsats"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Viktiga lärdomar från vår IaC-resa: Implementering av IaC kräver både teknisk excellens och organisatorisk förändring. Framgångsrika transformationer karaktäriseras av stark ledningsengagemang, omfatt..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtida utveckling och teknologiska trender: Cloud-native technologies, edge computing och artificiell intelligens driver nästa generation av Infrastructure as Code, som vi utforskade djupgående i [k..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Rekommendationer för organisationer: Baserat på vår genomgång från grundläggande principer till avancerade implementationer, bör organisationer påbörja sin IaC-journey med pilot projects som demonstre..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Slutord: Infrastructure as Code representerar mer än bara teknisk evolution - det är en fundamental förändring av hur vi tänker kring infrastruktur. Genom att embraca IaC-principer kan organisationer ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automatisering och CI/CD för slutsats"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadsoptimering inom slutsats"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för slutsats"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av slutsats"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för slutsats"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för slutsats"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Ordlista
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Ordlista"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Grundläggande koncept och verktyg: **API (Application Programming Interface):** Gränssnitt som möjliggör kommunikation mellan olika mjukvarukomponenter eller system genom standardiserade protokoll och..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Deployment och operationella koncept: **Blue-Green Deployment:** Deploymentstrategi där två identiska produktionsmiljöer (blå och grön) används för att möjliggöra snabb rollback och minimal downtime. ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadshantering och optimering: **FinOps:** Disciplin som kombinerar finansiell hantering med molnoperationer för att maximera affärsvärdet av molninvesteringar genom kostnadsoptimering och resource..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Testning och kvalitetssäkring: **Terratest:** Open source Go-bibliotek för automatiserad testning av Infrastructure as Code, särskilt designat för Terraform-moduler och cloud infrastructure. **Policy ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Strategiska och organisatoriska koncept: **Cloud-First Strategy:** Strategisk approach där organisationer primärt väljer molnbaserade lösningar för nya IT-initiativ innan on-premises alternativ övervä..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadsoptimering inom ordlista"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för ordlista"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av ordlista"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för ordlista"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för ordlista"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Om författarna
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Om författarna"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Gunnar Nordqvist: **Certifierad Chefsarkitekt och IT-säkerhetsspecialist** Gunnar Nordqvist är en erfaren IT-arkitekt med gedigen bakgrund inom sjukvårdssektorn och bred kompetens inom IT-tjänstehante..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Bidragsgivare och community: **Open Source Philosophy**: Denna bok är utvecklad med open source principles och community collaboration. Source code, examples och templates är tillgängliga för svenska ..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Svenska organisationers perspektiv på om författarna"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Säkerhetsaspekter inom om författarna"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automatisering och CI/CD för om författarna"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Kostnadsoptimering inom om författarna"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för om författarna"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av om författarna"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för om författarna"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för om författarna"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Appendix A: Kodexempel och tekniska implementationer
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Appendix A: Kodexempel och tekniska implementationer"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add diagram
    diagram_path = "../docs/images/diagram_26_appendix.png"
    if os.path.exists(diagram_path):
        try:
            slide.shapes.add_picture(diagram_path, Inches(0.5), Inches(1.2), Inches(4), Inches(3))
        except Exception as e:
            print(f"Warning: Could not add diagram {diagram_path}: {e}")
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Navigering i appendix: Kodexemplen är organiserade i följande kategorier: 1. **[CI/CD Pipelines och automatisering](#cicd-pipelines)** 2. **[Infrastructure as Code - Terraform](#terraform-iac)**"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• CI/CD Pipelines och automatisering {#cicd-pipelines}: Denna sektion innehåller alla exempel på CI/CD-pipelines, GitHub Actions workflows och automationsprocesser för svenska organisationer. *Refereras..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Infrastructure as Code - CloudFormation {#cloudformation-iac}: Denna sektion innehåller CloudFormation templates för AWS-infrastruktur anpassad för svenska organisationer. *Refereras från Kapitel 7: [..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Automation Scripts {#automation-scripts}: Denna sektion innehåller Python-skript och andra automationsverktyg för Infrastructure as Code-hantering. *Refereras från Kapitel 22: [Best practices och lärd..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Configuration Files {#configuration}: Denna sektion innehåller konfigurationsfiler för olika verktyg och tjänster. *Refereras från Kapitel 22: [Best practices och lärda läxor](22_best_practices.md)*"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Referenser och navigering: Varje kodexempel i denna appendix kan refereras från huvudtexten med dess unika identifierare. För att hitta specifika implementationer: 1. **Använd sökfunktion** - Sök efte..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Skalbarhet och prestanda för appendix a: kodexempel och tekniska implementationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Övervakning och loggning av appendix a: kodexempel och tekniska implementationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Team-organisation för appendix a: kodexempel och tekniska implementationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för appendix a: kodexempel och tekniska implementationer"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    # Chapter: Book Cover Design - \"Arkitektur som kod\"
    slide_layout = prs.slide_layouts[6]  # Blank layout for custom positioning
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Book Cover Design - \"Arkitektur som kod\""
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    # Add key points
    points_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4.5), Inches(6))
    points_frame = points_box.text_frame
    points_frame.text = "Viktiga punkter:"
    points_frame.paragraphs[0].font.size = Pt(16)
    points_frame.paragraphs[0].font.bold = True
    points_frame.paragraphs[0].font.color.rgb = RGBColor(0, 106, 167)  # Swedish blue
    

    p = points_frame.add_paragraph()
    p.text = "• Overview: Professional book cover design for Kvadrat\'s \"Arkitektur som kod\" (Infrastructure as Code) publication. The design follows Kvadrat\'s brand guidelines and incorporates modern visual eleme..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Design Variations: **File:** `templates/book-cover-final.html` - Modern gradient background (Kvadrat Blue to Dark Blue) - Advanced code architecture visual elements"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Export Formats: The design is available in multiple high-resolution formats: - **PDF**: 300 DPI print-ready format - **PNG**: High-resolution raster (2480×3508 pixels at 300 DPI)"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Technical Specifications: - **Format**: A4 (210mm × 297mm) - **Aspect Ratio**: Standard book cover proportions - **Resolution**: 300 DPI for print, 150 DPI for screen"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Usage Instructions: 1. Use `exports/book-cover/pdf/book-cover-print.pdf` 2. Ensure printer supports RGB color space (convert to CMYK if needed) 3. Recommended paper: High-quality matte or glossy finis..."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Generation Scripts: Generates all export formats automatically. Creates light and minimal design variations."
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Quality Assurance: - ✅ Kvadrat color palette strictly followed - ✅ Typography hierarchy maintained - ✅ Logo placement and sizing correct"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Maintenance: 1. Modify source HTML/CSS files in `templates/` 2. Run export generation script to update all formats 3. Test print quality and brand compliance"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Support: For questions about the design or technical implementation: - Review brand guidelines in `exports/book-cover/source/` - Check design system documentation"
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    p = points_frame.add_paragraph()
    p.text = "• Framtiden för book cover design - \"arkitektur som kod\""
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray

    
    # Save presentation
    output_path = "arkitektur_som_kod_presentation.pptx"
    prs.save(output_path)
    print(f"✅ Presentation saved to {output_path}")
    print(f"📊 Total slides created: {len(prs.slides)}")
    print("🎨 Styled with Swedish theme colors")
    print("📋 Each slide includes chapter title, diagram (when available), and 10 key points")

if __name__ == "__main__":
    create_presentation()
